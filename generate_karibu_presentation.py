from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = Path("Karibu_Grocery_Presentation_v2.pptx")


COLORS = {
    "forest": RGBColor(42, 92, 61),
    "deep_green": RGBColor(24, 59, 46),
    "sage": RGBColor(222, 233, 222),
    "mint": RGBColor(236, 244, 237),
    "sand": RGBColor(244, 235, 220),
    "gold": RGBColor(221, 161, 94),
    "navy": RGBColor(19, 36, 66),
    "slate": RGBColor(88, 104, 127),
    "soft_red": RGBColor(191, 67, 67),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(193, 203, 190),
}


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=20, color=None,
                bold=False, font_name="Aptos", align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["navy"]
    return shape


def add_round_box(slide, left, top, width, height, fill_color, line_color=None, radius_text=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    line = shape.line
    line.color.rgb = line_color or fill_color
    line.width = Pt(1)
    if radius_text:
        shape.text = radius_text
    return shape


def add_badge(slide, left, top, text, fill_color, text_color):
    badge = add_round_box(slide, left, top, Inches(1.55), Inches(0.34), fill_color, fill_color)
    frame = badge.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return badge


def add_bullets(slide, left, top, width, height, items, size=17, color=None):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.bullet = True
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["navy"]
    return shape


def add_header_band(slide, title, subtitle):
    add_round_box(slide, Inches(0.45), Inches(0.35), Inches(12.4), Inches(0.92), COLORS["sand"])
    add_badge(slide, Inches(0.72), Inches(0.56), "Karibu Grocery", COLORS["forest"], COLORS["white"])
    add_textbox(slide, Inches(2.45), Inches(0.54), Inches(6.4), Inches(0.32), title, size=24, bold=True)
    add_textbox(slide, Inches(2.45), Inches(0.83), Inches(7.8), Inches(0.24), subtitle, size=11, color=COLORS["slate"])


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["mint"])

    add_round_box(slide, Inches(0.55), Inches(0.55), Inches(6.15), Inches(5.85), COLORS["forest"])
    add_round_box(slide, Inches(6.95), Inches(0.55), Inches(5.85), Inches(5.85), COLORS["sand"])
    add_badge(slide, Inches(1.0), Inches(1.0), "Branch-Aware Retail Platform", COLORS["gold"], COLORS["deep_green"])
    add_textbox(slide, Inches(1.0), Inches(1.6), Inches(5.0), Inches(1.3),
                "Karibu Grocery\nManagement System", size=28, color=COLORS["white"], bold=True)
    add_textbox(slide, Inches(1.0), Inches(3.15), Inches(4.9), Inches(1.1),
                "A modern operations platform for sales, procurement, accounting, reporting, and secure user control.",
                size=18, color=COLORS["sage"])
    add_textbox(slide, Inches(1.0), Inches(5.45), Inches(4.8), Inches(0.3),
                "Pitch deck generated on March 11, 2026", size=11, color=COLORS["sage"])

    preview = add_round_box(slide, Inches(7.4), Inches(1.05), Inches(4.95), Inches(4.65), COLORS["white"], COLORS["line"])
    preview.line.width = Pt(1.5)
    add_textbox(slide, Inches(7.75), Inches(1.35), Inches(4.1), Inches(0.45), "Product Snapshot", size=18, bold=True)
    for i, label in enumerate(["Dashboard", "Procurement", "Accounting", "Reports", "User Admin"]):
        row = i % 3
        col = 0 if i < 3 else 1
        x = 7.75 + (2.1 * col)
        y = 2.0 + (1.05 * row)
        card = add_round_box(slide, Inches(x), Inches(y), Inches(1.8), Inches(0.82),
                             COLORS["mint"] if i % 2 == 0 else COLORS["sage"], COLORS["line"])
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.15), Inches(1.45), Inches(0.22), label, size=12, bold=True)
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.40), Inches(1.45), Inches(0.16), "Live module view", size=9, color=COLORS["slate"])
    add_textbox(slide, Inches(7.75), Inches(5.9), Inches(4.0), Inches(0.22),
                "Built for multi-branch retail and distribution teams", size=11, color=COLORS["slate"])


def create_problem_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_header_band(slide, "The Problem and the Karibu Solution",
                    "Replacing fragmented manual work with one connected operational workflow")

    add_round_box(slide, Inches(0.65), Inches(1.55), Inches(5.75), Inches(4.95), COLORS["sand"], COLORS["line"])
    add_textbox(slide, Inches(0.95), Inches(1.9), Inches(3.5), Inches(0.4), "Business Pain Points", size=22, bold=True)
    add_bullets(slide, Inches(0.95), Inches(2.45), Inches(4.9), Inches(3.5), [
        "Spreadsheets, paper records, and chat-based coordination create delays and errors.",
        "Stock visibility becomes weak across teams and branches.",
        "Reporting is slow, disconnected, and hard to trust.",
        "Accountability drops when tools and responsibilities are spread across many systems.",
    ])

    add_round_box(slide, Inches(6.7), Inches(1.55), Inches(5.9), Inches(4.95), COLORS["mint"], COLORS["line"])
    add_textbox(slide, Inches(7.0), Inches(1.9), Inches(3.9), Inches(0.4), "How Karibu Solves It", size=22, bold=True)
    add_bullets(slide, Inches(7.0), Inches(2.45), Inches(5.0), Inches(3.5), [
        "Centralizes sales, procurement, accounting, reporting, and administration.",
        "Supports branch-aware access with role-based permissions for every team.",
        "Captures daily transactions in one system for cleaner records and faster decisions.",
        "Turns the backend into a practical, user-friendly product for real operations.",
    ])


def create_platform_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["mint"])
    add_header_band(slide, "Platform Overview",
                    "One connected product experience across frontline work and management insight")

    modules = [
        ("Sales", "POS-style selling, receipts, payment tracking, and daily reporting."),
        ("Procurement", "Purchases, supplier flows, stock records, adjustments, and balance tracking."),
        ("Accounting", "Expenses, other income, and credit collection management."),
        ("Reports", "Cross-module reports with export-ready summaries for decision-making."),
        ("User Admin", "Role-based approvals, branch visibility, and access control."),
        ("Dashboard", "Branch health, sales activity, and stock alert visibility at a glance."),
    ]

    positions = [
        (0.75, 1.65), (4.35, 1.65), (7.95, 1.65),
        (0.75, 4.0), (4.35, 4.0), (7.95, 4.0),
    ]

    for index, (title, text) in enumerate(modules):
        x, y = positions[index]
        fill = COLORS["white"] if index % 2 == 0 else COLORS["sand"]
        add_round_box(slide, Inches(x), Inches(y), Inches(3.1), Inches(1.9), fill, COLORS["line"])
        add_textbox(slide, Inches(x + 0.18), Inches(y + 0.18), Inches(2.3), Inches(0.28), title, size=18, bold=True)
        add_textbox(slide, Inches(x + 0.18), Inches(y + 0.58), Inches(2.65), Inches(0.9), text, size=12, color=COLORS["slate"])
        add_badge(slide, Inches(x + 2.0), Inches(y + 1.42), "Core Module", COLORS["forest"], COLORS["white"])


def create_ui_walkthrough_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_header_band(slide, "Interface Walkthrough",
                    "A stylized walkthrough based on the shared UI screenshots")

    sections = [
        ("Dashboard", "Highlights low stock, today's sales, quick app entry points, and branch status."),
        ("Accounting", "Tracks expenses, other income, and credit collection with export support."),
        ("Reports", "Summarizes inflow, outflow, exposure, and module-level performance in one view."),
        ("Procurement", "Supports purchase entry, supplier selection, cost capture, and receipt printing."),
        ("Suppliers", "Maintains supplier contacts, terms, balances, and payment workflow."),
        ("User Console", "Shows approvals, role management, permissions, and branch overview."),
    ]

    positions = [
        (0.55, 1.7), (4.25, 1.7), (7.95, 1.7),
        (0.55, 4.2), (4.25, 4.2), (7.95, 4.2),
    ]

    for i, (title, desc) in enumerate(sections):
        x, y = positions[i]
        panel = add_round_box(slide, Inches(x), Inches(y), Inches(3.0), Inches(1.95),
                              COLORS["mint"], COLORS["line"])
        topbar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(3.0), Inches(0.34)
        )
        topbar.fill.solid()
        topbar.fill.fore_color.rgb = COLORS["sand"] if i % 2 == 0 else COLORS["sage"]
        topbar.line.color.rgb = COLORS["line"]
        add_textbox(slide, Inches(x + 0.16), Inches(y + 0.08), Inches(1.7), Inches(0.16), title, size=12, bold=True)

        for mini in range(3):
            mx = x + 0.18 + (mini * 0.93)
            my = y + 0.55
            add_round_box(slide, Inches(mx), Inches(my), Inches(0.72), Inches(0.42),
                          COLORS["white"], COLORS["line"])
        add_round_box(slide, Inches(x + 0.18), Inches(y + 1.08), Inches(2.62), Inches(0.34), COLORS["white"], COLORS["line"])
        add_round_box(slide, Inches(x + 0.18), Inches(y + 1.5), Inches(1.15), Inches(0.24), COLORS["forest"], COLORS["forest"])
        caption = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 1.72), Inches(2.72), Inches(0.36))
        frame = caption.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.TOP
        p = frame.paragraphs[0]
        p.text = desc
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(8.5)
        run.font.color.rgb = COLORS["slate"]


def create_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["mint"])
    add_header_band(slide, "Technology and Deployment",
                    "A modern web stack designed for maintainability, scale, and secure operations")

    left = add_round_box(slide, Inches(0.75), Inches(1.7), Inches(3.55), Inches(4.7), COLORS["white"], COLORS["line"])
    mid = add_round_box(slide, Inches(4.65), Inches(1.7), Inches(3.0), Inches(4.7), COLORS["sand"], COLORS["line"])
    right = add_round_box(slide, Inches(7.95), Inches(1.7), Inches(4.0), Inches(4.7), COLORS["white"], COLORS["line"])
    for box in (left, mid, right):
        box.line.width = Pt(1.5)

    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(2.0), Inches(0.3), "Frontend", size=22, bold=True)
    add_bullets(slide, Inches(1.0), Inches(2.45), Inches(2.85), Inches(3.2), [
        "Vue 3",
        "Vite",
        "Pinia",
        "Vue Router",
        "Tailwind CSS",
        "Axios",
        "Playwright",
    ], size=15)

    add_textbox(slide, Inches(4.95), Inches(2.0), Inches(2.0), Inches(0.3), "Backend", size=22, bold=True)
    add_bullets(slide, Inches(4.95), Inches(2.45), Inches(2.15), Inches(3.2), [
        "Node.js",
        "Express.js",
        "MongoDB + Mongoose",
        "JWT auth",
        "bcryptjs",
        "Swagger",
        "Bruno",
    ], size=15)

    add_textbox(slide, Inches(8.25), Inches(2.0), Inches(2.6), Inches(0.3), "Cloud Delivery", size=22, bold=True)
    add_bullets(slide, Inches(8.25), Inches(2.45), Inches(3.2), Inches(2.2), [
        "Vercel for frontend deployment",
        "Render for backend hosting",
        "MongoDB Atlas for cloud database operations",
        "Export-ready reporting and centralized branch visibility",
    ], size=14)

    connector = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(3.95), Inches(3.2), Inches(0.45), Inches(0.6)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = COLORS["forest"]
    connector.line.color.rgb = COLORS["forest"]
    connector2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(7.45), Inches(3.2), Inches(0.45), Inches(0.6)
    )
    connector2.fill.solid()
    connector2.fill.fore_color.rgb = COLORS["forest"]
    connector2.line.color.rgb = COLORS["forest"]


def create_value_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_header_band(slide, "Business Value",
                    "Why this platform matters to growing retail and distribution businesses")

    add_round_box(slide, Inches(0.8), Inches(1.8), Inches(2.7), Inches(3.95), COLORS["forest"], COLORS["forest"])
    add_textbox(slide, Inches(1.1), Inches(2.2), Inches(2.0), Inches(0.35), "Key Strengths", size=22, color=COLORS["white"], bold=True)
    add_bullets(slide, Inches(1.1), Inches(2.8), Inches(2.0), Inches(2.4), [
        "Multi-branch support",
        "Role-based access",
        "Operational frontend",
        "Transaction-safe workflows",
        "Clear reporting insight",
    ], size=14, color=COLORS["sage"])

    metric_cards = [
        ("Faster Decisions", "Management gets live branch-level visibility instead of waiting for manual reports."),
        ("Better Accountability", "Transactions, supplier records, and approvals live in one auditable workflow."),
        ("Operational Efficiency", "Teams sell, procure, report, and manage users without switching systems."),
        ("Scalable Growth", "The same platform supports more branches, more users, and cleaner oversight."),
    ]
    positions = [(3.9, 1.9), (8.1, 1.9), (3.9, 4.1), (8.1, 4.1)]
    fills = [COLORS["sand"], COLORS["mint"], COLORS["mint"], COLORS["sand"]]
    for i, (title, text) in enumerate(metric_cards):
        x, y = positions[i]
        add_round_box(slide, Inches(x), Inches(y), Inches(3.55), Inches(1.75), fills[i], COLORS["line"])
        add_textbox(slide, Inches(x + 0.18), Inches(y + 0.18), Inches(2.6), Inches(0.25), title, size=17, bold=True)
        add_textbox(slide, Inches(x + 0.18), Inches(y + 0.56), Inches(3.0), Inches(0.75), text, size=11, color=COLORS["slate"])


def create_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["forest"])

    add_round_box(slide, Inches(0.75), Inches(0.8), Inches(11.8), Inches(5.85), COLORS["deep_green"], COLORS["deep_green"])
    add_badge(slide, Inches(1.2), Inches(1.25), "Why Karibu Stands Out", COLORS["gold"], COLORS["deep_green"])
    add_textbox(slide, Inches(1.2), Inches(1.85), Inches(7.8), Inches(1.0),
                "One branch-aware platform for\nsales, stock, finance, reporting, and control.",
                size=28, color=COLORS["white"], bold=True)
    add_textbox(slide, Inches(1.2), Inches(3.2), Inches(8.4), Inches(1.0),
                "Karibu does more than digitize one process. It creates a complete operational system for both frontline staff and leadership.",
                size=18, color=COLORS["sage"])
    add_textbox(slide, Inches(1.2), Inches(5.55), Inches(3.3), Inches(0.3),
                "Prepared for presentation use", size=12, color=COLORS["sage"])

    quote_box = add_round_box(slide, Inches(8.85), Inches(1.55), Inches(2.85), Inches(3.65), COLORS["sand"], COLORS["sand"])
    quote_box.line.width = Pt(0)
    add_textbox(slide, Inches(9.15), Inches(2.0), Inches(2.2), Inches(0.35), "Value Proposition", size=18, bold=True)
    add_textbox(slide, Inches(9.15), Inches(2.55), Inches(2.1), Inches(1.7),
                "Karibu helps growing businesses move from fragmented manual operations to a connected, reliable workflow.",
                size=15, color=COLORS["navy"])
    add_textbox(slide, Inches(9.15), Inches(4.65), Inches(1.8), Inches(0.2), "karibu grocery system", size=10, color=COLORS["slate"])


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_problem_solution_slide(prs)
    create_platform_overview_slide(prs)
    create_ui_walkthrough_slide(prs)
    create_architecture_slide(prs)
    create_value_slide(prs)
    create_closing_slide(prs)

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_presentation()
